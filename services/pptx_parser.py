"""
PowerPoint (.pptx) parser.

Content documents only (training decks, SOPs-as-slides, etc.), ingested
exactly like PDFs/DOCX/Excel: extract_pages() turns the .pptx into
PageContent objects (the same dataclass PDFParser produces) so the file
flows through the existing DocumentChunker / SemanticChunkingService /
Qdrant pipeline unchanged.

Unlike DocxParser (which has no native notion of "pages" and returns a
single PageContent for the whole file), a .pptx has a natural per-slide
unit. extract_pages() returns one PageContent per slide, with
page_number/page_label set to the slide number -- so citations in the app
can say "Slide 4" the same way PDF citations say "Page 4".

Only text content is extracted (title/body placeholders, other text
boxes, and table cells). Slide images/diagrams are not OCR'd here --
route standalone photos of slides through services/image_parser.py
instead; a .pptx's embedded pictures are left as-is since captions/notes
usually carry the same information as text on the slide.
"""

import logging
import os
from typing import List

from pptx import Presentation

from services.pdf_parser import PageContent

logger = logging.getLogger(__name__)


class PptxParser:
    """Extracts per-slide text content from .pptx files."""

    def __init__(self, pptx_folder: str):
        self.pptx_folder = pptx_folder

    def _list_pptx_files(self) -> List[str]:
        return [
            os.path.join(self.pptx_folder, f)
            for f in os.listdir(self.pptx_folder)
            if f.lower().endswith(".pptx")
        ]

    @staticmethod
    def _slide_text(slide) -> str:
        lines: List[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if not text:
                        # python-pptx sometimes puts all text in para.text
                        # directly when there are no explicit runs.
                        text = para.text.strip()
                    if text:
                        lines.append(text)
            elif shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        lines.append(" | ".join(cells))

        # Speaker notes often carry real content (script/context) for
        # training decks -- include them, clearly separated.
        if slide.has_notes_slide:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes_text:
                lines.append(f"[Speaker notes: {notes_text}]")

        return "\n".join(lines)

    def extract_pages(self, file_path: str) -> List[PageContent]:
        """Read a .pptx as training content and return one PageContent per slide."""
        filename = os.path.basename(file_path)
        presentation = Presentation(file_path)

        core = presentation.core_properties
        metadata = {
            "title": core.title or "",
            "author": core.author or "",
            "toc_section": "",
        }

        pages: List[PageContent] = []
        for idx, slide in enumerate(presentation.slides, start=1):
            text = self._slide_text(slide)
            if not text.strip():
                continue
            pages.append(
                PageContent(
                    filename=filename,
                    page_number=idx,
                    page_label=f"Slide {idx}",
                    text=text,
                    metadata=dict(metadata),
                )
            )

        if not pages:
            logger.warning("No extractable text found in '%s'.", filename)

        return pages

    def extract_all(self) -> List[PageContent]:
        all_pages: List[PageContent] = []
        pptx_files = self._list_pptx_files()
        logger.info("Found %d .pptx file(s) in '%s'", len(pptx_files), self.pptx_folder)

        for path in pptx_files:
            try:
                extracted = self.extract_pages(path)
                logger.info("Extracted %d slide(s) from '%s'", len(extracted), os.path.basename(path))
                all_pages.extend(extracted)
            except Exception as exc:
                logger.error("Failed to process '%s': %s", path, exc)

        return all_pages